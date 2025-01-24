import os
import argparse
from pptx import Presentation
from pdf2image import convert_from_path
import cv2
import numpy as np

def process_pptx(pptx_file, output_folder):
    prs = Presentation(pptx_file)
    print(f"Opened presentation: {pptx_file}")
    
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    print(f"Output folder: {output_folder}")

    # Get slide dimensions in pixels (PPTX uses EMUs, 1 inch = 914400 EMU, 1 inch = 96 pixels)
    slide_width_emu = prs.slide_width
    slide_height_emu = prs.slide_height
    slide_width_px = int(slide_width_emu / 914400 * 96)
    slide_height_px = int(slide_height_emu / 914400 * 96)

    print(f"PPTX slide size in pixels: {slide_width_px}x{slide_height_px}")

    # Convert slides to images
    print("Converting slides to images...")
    slides_images = convert_from_path(pptx_file.replace('.pptx', '.pdf'), dpi=150)
    
    # Get converted image dimensions
    img_width, img_height = slides_images[0].size
    print(f"Converted image size: {img_width}x{img_height}")

    # Calculate scaling factors to align image with PPTX coordinates
    scale_x = img_width / slide_width_px
    scale_y = img_height / slide_height_px
    print(f"Scale factors - X: {scale_x}, Y: {scale_y}")

    for slide_idx, slide_img in enumerate(slides_images[5:], start=6):  # Start from slide 6
        print(f"Processing slide {slide_idx}")

        # Convert to OpenCV format
        slide_np = cv2.cvtColor(np.array(slide_img), cv2.COLOR_RGB2BGR)

        slide_title = ""
        for shape in prs.slides[slide_idx - 1].shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                slide_title = shape.text_frame.text.strip().replace(' ', '_')
                break

        print(f"Slide title: {slide_title if slide_title else 'Not Found'}")

        box_counter = 1
        for shape in prs.slides[slide_idx - 1].shapes:
            if shape.shape_type == 1 and shape.line and shape.line.color and hasattr(shape.line.color, 'rgb') and shape.line.color.rgb == (255, 0, 0):  # AUTOSHAPE with red border
                # Apply scaling factors to coordinates
                left = int((shape.left / 9525) * scale_x)
                top = int((shape.top / 9525) * scale_y)
                width = int((shape.width / 9525) * scale_x)
                height = int((shape.height / 9525) * scale_y)

                cropped_img = slide_np[top:top + height, left:left + width]
                
                # Assign sequential number for slide 6
                box_label = str(box_counter)
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        box_label = paragraph.text.strip()
                        break
                output_filename = f'{slide_title}s{box_label}.jpg'
                output_path = os.path.join(output_folder, output_filename)
                
                success = cv2.imwrite(output_path, cropped_img, [cv2.IMWRITE_JPEG_QUALITY, 85])
                if success:
                    print(f"Successfully saved: {output_path}")
                else:
                    print(f"Failed to save: {output_path}")

                box_counter += 1

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Extract red squares from PowerPoint slides.")
    parser.add_argument("pptx_file", type=str, default='resources/mario_scenes_manual_annotation.pptx', nargs='?', help="Path to the PowerPoint file")
    parser.add_argument("output_folder", type=str, default='resources/scenes_images', nargs='?', help="Path to the output folder")
    args = parser.parse_args()

    process_pptx(args.pptx_file, args.output_folder)
